"""Build two PowerPoint presentations for the Agentic Health Interop framework.

1. Executive deck (President) - 6-8 slides, vision + benefits + simple diagram
2. PM/Engineer deck - 8-10 slides, detailed use cases + architecture + framework vs content

Requires: pip install python-pptx Pillow
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DOCS = os.path.dirname(__file__)
IMG = os.path.join(DOCS, "img")

# --- Brand colors ---
DARK_BG    = RGBColor(0x0B, 0x0D, 0x11)
PANEL_BG   = RGBColor(0x13, 0x17, 0x1E)
CARD_BG    = RGBColor(0x1A, 0x1F, 0x2A)
ACCENT     = RGBColor(0x60, 0x6F, 0xF3)  # indigo accent
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
TEAL       = RGBColor(0x14, 0xB8, 0xA6)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
RED        = RGBColor(0xEF, 0x44, 0x44)
WHITE      = RGBColor(0xE6, 0xE8, 0xEB)
MUTED      = RGBColor(0x8B, 0x95, 0xA5)
LIGHT_BG   = RGBColor(0xF8, 0xF9, 0xFA)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
MID_TEXT    = RGBColor(0x4A, 0x4A, 0x5A)
BLUE_CORP   = RGBColor(0x00, 0x50, 0x8F)  # InterSystems-ish blue
BLUE_LIGHT  = RGBColor(0x00, 0x7B, 0xC0)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)

W = Inches(13.333)
H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=DARK_TEXT,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, size=18, color=DARK_TEXT, bold=False, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = space_before
    return p


def add_bullet(tf, text, size=16, color=DARK_TEXT, level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = level
    p.space_before = Pt(4)
    return p


def add_rect(slide, left, top, width, height, fill_color, text="",
             text_color=WHITE, text_size=12, bold=False, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(text_size)
        tf.paragraphs[0].font.color.rgb = text_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.name = "Calibri"
        shape.text_frame.margin_left = Pt(6)
        shape.text_frame.margin_right = Pt(6)
        shape.text_frame.margin_top = Pt(4)
        shape.text_frame.margin_bottom = Pt(4)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=MUTED):
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT = 1
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        else:
            slide.shapes.add_picture(path, left, top)


# ============================================================================
# DECK 1: EXECUTIVE (PRESIDENT) - 7 slides
# ============================================================================
def build_exec_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)
    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Agentic Health Interoperability Framework",
             size=36, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
             "An AI-native framework for building healthcare integration copilots on IRIS for Health",
             size=20, color=MUTED, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             "InterSystems AI Hub  |  May 2026",
             size=14, color=MUTED, align=PP_ALIGN.LEFT)
    # accent line
    add_rect(slide, Inches(1), Inches(2.55), Inches(3), Pt(3), ACCENT)

    # --- Slide 2: The Opportunity ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "The Opportunity", size=30, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2), Pt(3), ACCENT)

    # Problem column
    add_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), RGBColor(0xFF,0xF1,0xF2), radius=True)
    tf = add_text(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5),
             "The Problem", size=20, color=RED, bold=True)
    add_bullet(tf, "Integration engineers spend 70% of their time on repetitive tasks: wiring productions, writing DTLs, debugging HL7 segment mappings", size=14, color=DARK_TEXT)
    add_bullet(tf, "Deep IRIS expertise is scarce -- new engineers face a steep learning curve on ObjectScript, Productions, DTL, and the HS.* class library", size=14, color=DARK_TEXT)
    add_bullet(tf, "Every customer builds the same plumbing from scratch: ADT routing, HL7-to-FHIR pipelines, error triage dashboards", size=14, color=DARK_TEXT)
    add_bullet(tf, "No AI assistance exists today for IRIS-specific interoperability tasks -- general-purpose LLMs don't know Productions, SDA3, or the Ens.* class catalog", size=14, color=DARK_TEXT)

    # Solution column
    add_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(0xF0,0xFD,0xF4), radius=True)
    tf = add_text(slide, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.5),
             "The Solution", size=20, color=GREEN, bold=True)
    add_bullet(tf, "An AI framework built into IRIS for Health that understands Productions, DTL, HL7, FHIR, and the complete class library", size=14, color=DARK_TEXT)
    add_bullet(tf, "Natural-language interface: engineers describe what they want, the agent builds it -- searching catalogs, proposing configurations, asking for approval, and validating results", size=14, color=DARK_TEXT)
    add_bullet(tf, "Out-of-the-box agents, tools, and skills that cover the most common integration patterns -- ready to use on day one", size=14, color=DARK_TEXT)
    add_bullet(tf, "Extensible framework: customers, SEs, and the dev community can add their own agents, MCPs, tools, and skills without modifying the core", size=14, color=DARK_TEXT)

    # --- Slide 3: Vision - What is the Framework ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Vision: An AI-Native Interoperability Platform", size=30, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2.5), Pt(3), ACCENT)

    add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8),
             "The %AI Framework provides the primitives. We build the healthcare layer on top.",
             size=16, color=MID_TEXT)

    # Simple layered diagram
    # Layer 1: Integration Engineer (top)
    add_rect(slide, Inches(2.5), Inches(2.2), Inches(8), Inches(0.7),
             RGBColor(0xE0,0xE7,0xFF), "Integration Engineer (natural language)",
             text_color=DARK_TEXT, text_size=14, bold=True, radius=True)

    # Layer 2: Chat UI + Admin UI
    add_rect(slide, Inches(2.5), Inches(3.15), Inches(3.8), Inches(0.6),
             ACCENT, "Chat UI (streaming)",
             text_color=WHITE, text_size=12, bold=True, radius=True)
    add_rect(slide, Inches(6.7), Inches(3.15), Inches(3.8), Inches(0.6),
             PURPLE, "Admin UI (configuration)",
             text_color=WHITE, text_size=12, bold=True, radius=True)

    # Layer 3: Agent + Skills
    add_rect(slide, Inches(2.5), Inches(4.0), Inches(8), Inches(0.6),
             BLUE_CORP, "Health Interop Agent  +  12 Domain Skills  +  Policies",
             text_color=WHITE, text_size=13, bold=True, radius=True)

    # Layer 4: MCP Servers
    y4 = Inches(4.85)
    w4 = Inches(1.85)
    for i, (name, clr) in enumerate([
        ("Production", TEAL), ("Transform", GREEN),
        ("Testing", AMBER), ("Catalog", BLUE_LIGHT)]):
        add_rect(slide, Inches(2.5) + i*(w4+Inches(0.2)), y4, w4, Inches(0.55),
                 clr, f"MCP: {name}", text_color=WHITE, text_size=11, bold=True, radius=True)

    # Layer 5: 42 Tools
    add_rect(slide, Inches(2.5), Inches(5.65), Inches(8), Inches(0.55),
             CARD_BG, "42 Tools  (ObjectScript + SQL + Embedded Python)",
             text_color=WHITE, text_size=12, bold=False, radius=True)

    # Layer 6: IRIS
    add_rect(slide, Inches(2.5), Inches(6.45), Inches(8), Inches(0.6),
             DARK_BG, "IRIS for Health  |  Vector Search  |  Secured Wallet  |  Productions  |  %Dictionary",
             text_color=MUTED, text_size=11, bold=False, radius=True)

    # --- Slide 4: Three Use Cases ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Three Core Use Cases", size=30, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2), Pt(3), ACCENT)

    # Three cards side by side
    card_w = Inches(3.7)
    card_h = Inches(5.0)
    gap = Inches(0.35)
    x0 = Inches(0.8)
    y0 = Inches(1.5)

    for i, (title, color, icon, items) in enumerate([
        ("Build Productions", TEAL, "BUILD",
         ["Describe what you need in plain English",
          "Agent searches class catalog for right Business Hosts",
          "Proposes production layout for approval",
          "Creates production, adds hosts, configures settings",
          "Sends test messages and validates end-to-end"]),
        ("Review & Improve", AMBER, "OPERATE",
         ["Triage errors across all productions",
          "Identify bottlenecks and recommend settings",
          "Review DTLs for hardcoded values and missing checks",
          "Suggest modernization using newer IRIS features",
          "Verify fixes with automated validation"]),
        ("Create Transformations", GREEN, "TRANSFORM",
         ["Trace data flow: HL7 v2 -> SDA3 -> FHIR R4",
          "Introspect HL7 schemas at sub-field level",
          "Create DTLs with correct source/target classes",
          "Dry-run transformations against sample data",
          "Field-level gap analysis via the Transformation and Mapping Catalog"]),
    ]):
        x = x0 + i * (card_w + gap)
        add_rect(slide, x, y0, card_w, card_h, RGBColor(0xFF,0xFF,0xFF), radius=True)
        # Color bar at top of card
        add_rect(slide, x, y0, card_w, Pt(4), color)
        # Category tag
        add_rect(slide, x + Inches(0.3), y0 + Inches(0.3), Inches(1.3), Inches(0.35),
                 color, icon, text_color=WHITE, text_size=10, bold=True, radius=True)
        # Title
        add_text(slide, x + Inches(0.3), y0 + Inches(0.8), card_w - Inches(0.6), Inches(0.5),
                 title, size=18, color=DARK_TEXT, bold=True)
        # Bullets
        tf = add_text(slide, x + Inches(0.3), y0 + Inches(1.3), card_w - Inches(0.6), Inches(3.5),
                      "", size=13, color=MID_TEXT)
        for item in items:
            add_bullet(tf, item, size=12, color=MID_TEXT)

    # --- Slide 5: Capabilities - Out of the Box vs Community ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Capabilities: Framework + Content", size=30, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2.5), Pt(3), ACCENT)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.6),
             "The framework ships with production-ready content. Customers and the community extend it without modifying the core.",
             size=15, color=MID_TEXT)

    # Two columns
    col_w = Inches(5.6)
    col_h = Inches(4.8)

    # Left: Out of the Box
    add_rect(slide, Inches(0.8), Inches(1.9), col_w, col_h, RGBColor(0xEF,0xF6,0xFF), radius=True)
    add_rect(slide, Inches(0.8), Inches(1.9), col_w, Pt(4), BLUE_CORP)
    add_text(slide, Inches(1.1), Inches(2.1), Inches(4), Inches(0.4),
             "Out of the Box (shipped by InterSystems)", size=16, color=BLUE_CORP, bold=True)
    tf = add_text(slide, Inches(1.1), Inches(2.65), col_w - Inches(0.6), Inches(3.8),
                  "", size=13, color=DARK_TEXT)
    add_bullet(tf, "1 Agent: Health Interop (router with orchestration policies)", size=13, color=DARK_TEXT, bold=True)
    add_bullet(tf, "4 MCP Servers: Production, Transform, Testing, Catalog", size=13, color=DARK_TEXT, bold=True)
    add_bullet(tf, "5 ToolSets with 42 Tools covering the full integration lifecycle", size=13, color=DARK_TEXT, bold=True)
    add_bullet(tf, "12 Skills: Productions, DTL, BPL, HL7 v2, FHIR R4, SDA, CDA, X12, Routing Rules, ESB Patterns, Adapters, REST", size=13, color=DARK_TEXT, bold=True)
    add_bullet(tf, "2 Vector Catalogs: Ens.* (164 classes) + HS.* (58 classes)", size=13, color=DARK_TEXT, bold=True)
    add_bullet(tf, "Transformation and Mapping Catalog: 1,538 field-level mappings across HL7/SDA3/FHIR", size=13, color=DARK_TEXT, bold=True)
    add_bullet(tf, "Admin UI + Chat UI, audit trail, connection manager", size=13, color=DARK_TEXT, bold=True)

    # Right: Community / Customer Content
    add_rect(slide, Inches(6.7), Inches(1.9), col_w, col_h, RGBColor(0xF5,0xF3,0xFF), radius=True)
    add_rect(slide, Inches(6.7), Inches(1.9), col_w, Pt(4), PURPLE)
    add_text(slide, Inches(7.0), Inches(2.1), Inches(4.5), Inches(0.4),
             "Community & Customer Content", size=16, color=PURPLE, bold=True)
    tf = add_text(slide, Inches(7.0), Inches(2.65), col_w - Inches(0.6), Inches(3.8),
                  "", size=13, color=DARK_TEXT)
    add_bullet(tf, "Custom Agents: specialty copilots (Lab, Pharmacy, Claims, Imaging)", size=13, color=DARK_TEXT)
    add_bullet(tf, "Custom MCP Servers: wrap customer-specific APIs, EHR integrations, payer endpoints", size=13, color=DARK_TEXT)
    add_bullet(tf, "Custom Tools: call internal REST services, query custom tables, invoke Python ML models", size=13, color=DARK_TEXT)
    add_bullet(tf, "Custom Skills: organization-specific knowledge (local HL7 profiles, DTL conventions, naming standards)", size=13, color=DARK_TEXT)
    add_bullet(tf, "SE-built accelerators: pre-built productions for common patterns (ADT routing, HL7-to-FHIR, claims processing)", size=13, color=DARK_TEXT)
    add_bullet(tf, "ISV extensions: third-party tool packs distributed via IPM / Open Exchange", size=13, color=DARK_TEXT)

    # --- Slide 6: The Use Case - Healthcare Agent Interop ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Use Case: Healthcare Agent Interop", size=30, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2.5), Pt(3), ACCENT)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.6),
             "What we built to prove the framework -- a working Copilot on IRIS for Health 2026.2",
             size=15, color=MID_TEXT)

    # Two screenshots side by side
    add_image_safe(slide, os.path.join(IMG, "15_chatbot.png"),
                   Inches(0.8), Inches(1.9), width=Inches(6))
    add_image_safe(slide, os.path.join(IMG, "13_transforms_hl7_fhir.png"),
                   Inches(7.0), Inches(1.9), width=Inches(5.8))

    add_text(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(0.4),
             "Chat: natural-language production builder with streaming + tool cards",
             size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.0), Inches(5.5), Inches(5.8), Inches(0.4),
             "Transformation and Mapping Catalog: field-level HL7 v2 -> SDA3 -> FHIR R4 mapping explorer",
             size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)

    # Stats bar
    stats_y = Inches(6.2)
    stats = [
        ("61", "ObjectScript Classes"),
        ("42", "AI Tools"),
        ("12", "Domain Skills"),
        ("4", "MCP Servers"),
        ("1,538", "Field Mappings"),
        ("222", "Indexed Classes"),
    ]
    stat_w = Inches(1.8)
    for i, (num, label) in enumerate(stats):
        x = Inches(0.8) + i * (stat_w + Inches(0.25))
        add_text(slide, x, stats_y, stat_w, Inches(0.4),
                 num, size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x, stats_y + Inches(0.35), stat_w, Inches(0.3),
                 label, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # --- Slide 7: Benefits (visual cards) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Why This Matters", size=32, color=WHITE, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2), Pt(3), ACCENT)

    # 3x2 card grid
    benefits = [
        # (icon_text, accent_color, title, description, metric)
        ("DAYS\n  MIN",   TEAL,
         "Faster Time-to-Value",
         "Engineers describe what they need in plain English. The agent builds, configures, and validates -- reducing production setup from days to minutes.",
         "10x faster"),
        ("DAY 1",   GREEN,
         "Lower Barrier to Entry",
         "New engineers are productive immediately. Natural-language guidance replaces the steep learning curve of ObjectScript, DTL, and the Ens.* class library.",
         "Zero ramp-up"),
        ("0\nERRORS",   ACCENT,
         "Higher Quality",
         "Automated PostBuildValidation catches misconfigurations, missing settings, and broken routing before they ever reach a live production environment.",
         "Built-in validation"),
        ("IRIS",   BLUE_LIGHT,
         "Platform Stickiness",
         "AI assistance embedded in the integration engine is a compelling reason to choose and stay on IRIS for Health. No competitor offers this today.",
         "Unique to IRIS"),
        ("1+1=3",   PURPLE,
         "Community Flywheel",
         "Every custom tool, skill, and agent built by customers, SEs, and ISVs benefits the entire ecosystem. Content compounds; the platform gets smarter over time.",
         "Network effects"),
        ("ONLY\nONE",   AMBER,
         "Market Differentiator",
         "No other integration platform ships an AI framework with domain-specific healthcare tools. This is a first-mover advantage in the agentic AI era.",
         "First in market"),
    ]

    card_w = Inches(3.7)
    card_h = Inches(2.55)
    gap_x = Inches(0.35)
    gap_y = Inches(0.35)
    x0 = Inches(0.8)
    y0 = Inches(1.4)

    for idx, (icon_text, accent, title, desc, metric) in enumerate(benefits):
        col = idx % 3
        row = idx // 3
        cx = x0 + col * (card_w + gap_x)
        cy = y0 + row * (card_h + gap_y)

        # Card background
        add_rect(slide, cx, cy, card_w, card_h, PANEL_BG, radius=True)

        # Accent stripe on left
        add_rect(slide, cx, cy, Pt(5), card_h, accent)

        # Icon / metric badge (top-right of card)
        badge_w = Inches(0.9)
        badge_h = Inches(0.55)
        badge = add_rect(slide, cx + card_w - badge_w - Inches(0.15), cy + Inches(0.15),
                         badge_w, badge_h, accent, icon_text,
                         text_color=WHITE, text_size=9, bold=True, radius=True)

        # Title
        add_text(slide, cx + Inches(0.25), cy + Inches(0.15), card_w - Inches(1.4), Inches(0.4),
                 title, size=16, color=WHITE, bold=True)

        # Description
        add_text(slide, cx + Inches(0.25), cy + Inches(0.6), card_w - Inches(0.5), Inches(1.3),
                 desc, size=11, color=MUTED)

        # Bottom metric tag
        add_rect(slide, cx + Inches(0.25), cy + card_h - Inches(0.5),
                 Inches(1.5), Inches(0.3), accent, metric,
                 text_color=WHITE, text_size=10, bold=True, radius=True)

    # --- Slide 8: Next Steps ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Next Steps", size=32, color=WHITE, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2), Pt(3), ACCENT)

    # Timeline-style layout: 3 phases
    phase_data = [
        ("Immediate", TEAL, "Fix & Ship", [
            "Fix 2 blocker bugs (Skill $ZF, Bedrock hang)",
            "Ship ToolFilter as framework default",
            "Document %AI.ToolMgr.ExecuteTool path",
        ]),
        ("Q3-Q4 2026", AMBER, "Field Trial", [
            "Package as IPM for 3 customer pilots",
            "Build 2-3 specialty agents (Lab, Pharmacy, Claims)",
            "Open tool/skill authoring to SE community",
        ]),
        ("2027", PURPLE, "Ecosystem", [
            "Publish to Open Exchange",
            "SE accelerator library (20+ patterns)",
            "Community marketplace for tools and skills",
        ]),
    ]

    phase_w = Inches(3.7)
    phase_h = Inches(3.8)
    phase_y = Inches(1.8)
    for i, (when, clr, label, items) in enumerate(phase_data):
        px = Inches(0.8) + i * (phase_w + Inches(0.35))

        # Phase card
        add_rect(slide, px, phase_y, phase_w, phase_h, PANEL_BG, radius=True)
        add_rect(slide, px, phase_y, phase_w, Pt(5), clr)

        # Phase number circle
        add_rect(slide, px + Inches(0.25), phase_y + Inches(0.25),
                 Inches(0.5), Inches(0.5), clr, str(i + 1),
                 text_color=WHITE, text_size=16, bold=True, radius=True)

        # When label
        add_text(slide, px + Inches(0.9), phase_y + Inches(0.25), Inches(2.5), Inches(0.35),
                 when, size=14, color=clr, bold=True)

        # What label
        add_text(slide, px + Inches(0.9), phase_y + Inches(0.55), Inches(2.5), Inches(0.3),
                 label, size=18, color=WHITE, bold=True)

        # Items
        tf = add_text(slide, px + Inches(0.25), phase_y + Inches(1.1),
                      phase_w - Inches(0.5), Inches(2.5), "", size=13, color=WHITE)
        for item in items:
            add_bullet(tf, item, size=13, color=MUTED)

    # Bottom bar
    add_text(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
             "Built on IRIS for Health 2026.2  |  %AI Framework build 162.0  |  61 classes, 42 tools, 12 skills  |  Ready for field trial",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

    path = os.path.join(DOCS, "Agentic_Health_Interop_Executive.pptx")
    prs.save(path)
    print(f"  saved {path}")


# ============================================================================
# DECK 2: PM / ENGINEER - 10 slides
# ============================================================================
def build_eng_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             "Agentic Health Interoperability Framework",
             size=36, color=WHITE, bold=True)
    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "Technical Deep-Dive: Architecture, Use Cases, and Extensibility",
             size=20, color=MUTED)
    add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "InterSystems AI Hub  |  May 2026  |  PM / Engineering Audience",
             size=14, color=MUTED)
    add_rect(slide, Inches(1), Inches(2.3), Inches(3), Pt(3), ACCENT)

    # --- Slide 2: The %AI Framework Primitives ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "The %AI Framework: Five Primitives", size=30, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2.5), Pt(3), ACCENT)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.6),
             "IRIS for Health 2026.2 ships five AI primitives. The framework provides the plumbing; the content makes it useful.",
             size=15, color=MID_TEXT)

    # Five boxes
    prims = [
        ("Agent", "%AI.Agent", "Orchestrates multi-step tasks.\nRoutes to tools, skills, policies.\nManages conversation context.", BLUE_CORP),
        ("MCP Server", "%AI.MCP.Service", "Groups related tools into\nnamed service domains.\nIn-process or HTTP/SSE transport.", TEAL),
        ("ToolSet", "%AI.ToolSet", "Collection of related tools.\nAuto-discovered by framework.\nJSON Schema input/output.", GREEN),
        ("Tool", "%AI.Tool", "Single capability callable by LLM.\nObjectScript, SQL, or Python.\nConfirmation gate support.", AMBER),
        ("Skill", "%AI.Agent.Skill", "Domain knowledge as text.\nInjected as sub-agent tool.\nUp to 32K chars of instructions.", PURPLE),
    ]
    box_w = Inches(2.2)
    box_h = Inches(3.2)
    gap = Inches(0.25)
    x0 = Inches(0.8)
    y0 = Inches(1.9)
    for i, (name, cls, desc, clr) in enumerate(prims):
        x = x0 + i * (box_w + gap)
        add_rect(slide, x, y0, box_w, box_h, RGBColor(0xFF,0xFF,0xFF), radius=True)
        add_rect(slide, x, y0, box_w, Pt(4), clr)
        add_text(slide, x + Inches(0.15), y0 + Inches(0.2), box_w - Inches(0.3), Inches(0.4),
                 name, size=16, color=DARK_TEXT, bold=True)
        add_text(slide, x + Inches(0.15), y0 + Inches(0.6), box_w - Inches(0.3), Inches(0.4),
                 cls, size=10, color=MUTED)
        add_text(slide, x + Inches(0.15), y0 + Inches(1.0), box_w - Inches(0.3), Inches(2),
                 desc, size=11, color=MID_TEXT)

    # Bottom: relationship
    add_text(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.6),
             "Binding chain:   Agent (1) --> MCP Servers (N) --> ToolSets (N) --> Tools (N)     |     Skills (N) <-> Agent (N)",
             size=14, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

    # Label: Framework vs Content
    add_rect(slide, Inches(0.8), Inches(6.1), Inches(5.5), Inches(0.6),
             BLUE_CORP, "FRAMEWORK (shipped by InterSystems, not modified)",
             text_color=WHITE, text_size=12, bold=True, radius=True)
    add_rect(slide, Inches(6.7), Inches(6.1), Inches(5.8), Inches(0.6),
             PURPLE, "CONTENT (agents, tools, skills -- built on top of the framework)",
             text_color=WHITE, text_size=12, bold=True, radius=True)

    # --- Slide 3: Architecture Diagram ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Architecture: Healthcare Agent Interop", size=30, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2.5), Pt(3), ACCENT)

    # Detailed layered diagram
    lm = Inches(1.5)
    full_w = Inches(10.3)

    # Row 1: UIs
    add_rect(slide, lm, Inches(1.5), Inches(4.9), Inches(0.65),
             ACCENT, "Chat UI (SSE streaming, tool-call cards, confirmation gate)",
             text_color=WHITE, text_size=11, bold=True, radius=True)
    add_rect(slide, lm + Inches(5.2), Inches(1.5), Inches(5.1), Inches(0.65),
             PURPLE, "Admin UI (agents, MCPs, tools, skills, connections, catalogs, transforms, audit)",
             text_color=WHITE, text_size=11, bold=True, radius=True)

    # Row 2: REST API
    add_rect(slide, lm, Inches(2.35), full_w, Inches(0.5),
             RGBColor(0xDD, 0xDD, 0xDD), "REST API: /api/agentic/*  (14 service classes, JWT/Basic auth, audit logging)",
             text_color=DARK_TEXT, text_size=11, radius=True)

    # Row 3: Agent + Policies + Skills
    add_rect(slide, lm, Inches(3.05), Inches(4.5), Inches(0.85),
             BLUE_CORP, "HealthInterop Agent\nManager + Monitor + SkillLoader",
             text_color=WHITE, text_size=11, bold=True, radius=True)
    add_rect(slide, lm + Inches(4.7), Inches(3.05), Inches(2.6), Inches(0.85),
             RGBColor(0x9B, 0x59, 0xB6), "Policies\nConfirmationGate\nToolFilter",
             text_color=WHITE, text_size=10, bold=True, radius=True)
    add_rect(slide, lm + Inches(7.5), Inches(3.05), Inches(2.8), Inches(0.85),
             RGBColor(0x8E, 0x44, 0xAD), "12 Skills\nHL7 | FHIR | DTL | BPL\nSDA | CDA | X12 | ...",
             text_color=WHITE, text_size=10, bold=True, radius=True)

    # Row 4: MCP Servers
    y4 = Inches(4.1)
    mcps = [
        ("Production\n13 tools", TEAL),
        ("Transform\n13 tools", GREEN),
        ("Testing\n8 tools", AMBER),
        ("Catalog\n8 tools", BLUE_LIGHT),
    ]
    mcp_w = Inches(2.4)
    for i, (label, clr) in enumerate(mcps):
        add_rect(slide, lm + i*(mcp_w + Inches(0.17)), y4, mcp_w, Inches(0.7),
                 clr, label, text_color=WHITE, text_size=10, bold=True, radius=True)

    # Row 5: Tool implementations
    add_rect(slide, lm, Inches(5.0), full_w, Inches(0.5),
             CARD_BG, "42 Tools: ObjectScript class methods  |  SQL queries  |  Embedded Python  |  REST wrappers",
             text_color=WHITE, text_size=11, radius=True)

    # Row 6: IRIS stores
    y6 = Inches(5.7)
    stores = [
        ("Vector Search\nHNSW + FastEmbed", DARK_BG),
        ("Secured Wallet\nAPI keys only", DARK_BG),
        ("Productions\nlive runtime", DARK_BG),
        ("%Dictionary\nclass metadata", DARK_BG),
    ]
    st_w = Inches(2.4)
    for i, (label, clr) in enumerate(stores):
        add_rect(slide, lm + i*(st_w + Inches(0.17)), y6, st_w, Inches(0.65),
                 clr, label, text_color=MUTED, text_size=10, radius=True)

    # Left labels
    labels = [
        (Inches(1.5), "UI Layer"),
        (Inches(2.35), "API Layer"),
        (Inches(3.05), "Agent Layer"),
        (Inches(4.1), "MCP Layer"),
        (Inches(5.0), "Tool Layer"),
        (Inches(5.7), "IRIS Stores"),
    ]
    for y, label in labels:
        add_text(slide, Inches(0.15), y, Inches(1.3), Inches(0.5),
                 label, size=9, color=MUTED, bold=True)

    # --- Slide 4: Use Case 1 - Build Productions ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Use Case 1: Build Productions", size=30, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2.5), Pt(3), TEAL)

    # Flow: 5 phases
    phases = [
        ("1. Discover", "Agent searches\nEns.* catalog for\nBusiness Hosts", TEAL),
        ("2. Propose", "Agent presents\nproduction layout\nfor approval", BLUE_CORP),
        ("3. Build", "Creates production,\nadds hosts, settings,\nrouting rules", GREEN),
        ("4. Validate", "PostBuildValidation\n+ test HL7 message\nthrough pipeline", AMBER),
        ("5. Report", "Summarizes what\nwas built, surfaces\nany warnings", PURPLE),
    ]
    ph_w = Inches(2.1)
    ph_h = Inches(1.6)
    for i, (title, desc, clr) in enumerate(phases):
        x = Inches(0.8) + i * (ph_w + Inches(0.25))
        add_rect(slide, x, Inches(1.4), ph_w, ph_h, RGBColor(0xFF,0xFF,0xFF), radius=True)
        add_rect(slide, x, Inches(1.4), ph_w, Pt(4), clr)
        add_text(slide, x + Inches(0.15), Inches(1.55), ph_w - Inches(0.3), Inches(0.35),
                 title, size=14, color=clr, bold=True)
        add_text(slide, x + Inches(0.15), Inches(1.95), ph_w - Inches(0.3), Inches(1),
                 desc, size=11, color=MID_TEXT)

    # Tools used
    tf = add_text(slide, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.4),
                  "Tools: search_ens, describe_class, create_production, add_business_host, update_settings, create_routing_rule, start/stop_production, PostBuildValidation, BuildAndSendHL7TestMessage",
                  size=11, color=MUTED)

    # Screenshot
    add_image_safe(slide, os.path.join(IMG, "15_chatbot.png"),
                   Inches(0.8), Inches(3.9), width=Inches(11.5))

    # --- Slide 5: Use Case 2 - Review & Improve ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Use Case 2: Review & Improve Existing Productions", size=30, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(3), Pt(3), AMBER)

    # Four capability cards
    caps = [
        ("Error Triage", "Query Event Log for recent errors.\nGroup by Business Host.\nIdentify suspended messages.\nRecommend remediation.", AMBER),
        ("Health Assessment", "Inspect production config.\nCheck queue depths and throughput.\nIdentify bottlenecks.\nRecommend settings changes.", TEAL),
        ("DTL Review", "Spot hardcoded values.\nFind missing null checks.\nFlag repeating-field bugs.\nSuggest refactored versions.", GREEN),
        ("Modernization", "Recommend newer IRIS features.\nReplace custom BPL with DTL.\nAdopt HL7-to-SDA-to-FHIR.\nUse record maps over parsers.", PURPLE),
    ]
    cap_w = Inches(2.8)
    cap_h = Inches(2.6)
    for i, (title, desc, clr) in enumerate(caps):
        x = Inches(0.8) + i * (cap_w + Inches(0.2))
        add_rect(slide, x, Inches(1.5), cap_w, cap_h, RGBColor(0xFF,0xFF,0xFF), radius=True)
        add_rect(slide, x, Inches(1.5), cap_w, Pt(4), clr)
        add_text(slide, x + Inches(0.2), Inches(1.7), cap_w - Inches(0.4), Inches(0.35),
                 title, size=15, color=clr, bold=True)
        add_text(slide, x + Inches(0.2), Inches(2.1), cap_w - Inches(0.4), Inches(2),
                 desc, size=11, color=MID_TEXT)

    tf = add_text(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.4),
                  "Tools: get_production, query_event_log, top_errors, query_message_status, message_summary, queue_status, list_dtls, describe_class",
                  size=11, color=MUTED)
    tf = add_text(slide, Inches(0.8), Inches(4.65), Inches(11.5), Inches(0.4),
                  "Skills: Productions, DTL, BPL, Adapters, ESBPattern -- deep domain knowledge for context-aware recommendations",
                  size=11, color=MUTED)

    # Screenshot of audit
    add_image_safe(slide, os.path.join(IMG, "14_audit.png"),
                   Inches(0.8), Inches(5.1), width=Inches(11.5))

    # --- Slide 6: Use Case 3 - Create Transformations ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Use Case 3: Create & Optimize Transformations", size=30, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(3), Pt(3), GREEN)

    # Data flow diagram
    add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
             "SDA3 is the universal pivot -- all external formats map through it",
             size=15, color=MID_TEXT)

    y_flow = Inches(2.0)
    add_rect(slide, Inches(1.5), y_flow, Inches(2.5), Inches(0.7),
             TEAL, "HL7 v2 (inbound)\nPID.11.3 City", text_color=WHITE, text_size=12, bold=True, radius=True)
    add_rect(slide, Inches(5.0), y_flow, Inches(3), Inches(0.7),
             BLUE_CORP, "SDA3 (canonical)\nAddress.City", text_color=WHITE, text_size=12, bold=True, radius=True)
    add_rect(slide, Inches(9.0), y_flow, Inches(2.8), Inches(0.7),
             GREEN, "FHIR R4 (outbound)\nAddress.city", text_color=WHITE, text_size=12, bold=True, radius=True)

    # Class names under the arrows
    add_text(slide, Inches(3.5), Inches(2.75), Inches(2), Inches(0.4),
             "HS.Gateway.HL7.\nHL7ToSDA3", size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.5), Inches(2.75), Inches(2), Inches(0.4),
             "HS.FHIR.DTL.SDA3.\nvR4.Address.Address", size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # Transformation and Mapping Catalog screenshot
    add_image_safe(slide, os.path.join(IMG, "13_transforms_hl7_fhir.png"),
                   Inches(0.8), Inches(3.4), width=Inches(11.5))

    # --- Slide 7: Framework vs Content ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Framework vs. Content: What Ships vs. What's Built", size=28, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(3), Pt(3), ACCENT)

    # Three-column layout: Framework | OOB Content | Custom Content
    col_w = Inches(3.7)
    col_h = Inches(5.2)
    cols = [
        ("IRIS %AI Framework", BLUE_CORP, "Shipped by InterSystems\nDo NOT modify", [
            "%AI.Agent runtime",
            "%AI.MCP.Service transport",
            "%AI.ToolSet / %AI.Tool",
            "%AI.Agent.Skill engine",
            "%AI.RAG.KnowledgeBase",
            "Rust LLM bridge (multi-provider)",
            "FastEmbed (384-dim vectors)",
            "JSON Schema validation",
        ]),
        ("Out-of-the-Box Content", TEAL, "Shipped as IPM package\nUpgradeable, overridable", [
            "HealthInterop Agent + policies",
            "4 MCP Servers",
            "42 Tools across 5 ToolSets",
            "12 Skills (HL7, FHIR, DTL, ...)",
            "2 Vector Catalogs (222 classes)",
            "Transformation and Mapping Catalog (1,538 field mappings)",
            "Admin UI + Chat UI",
            "Overlay pattern for customization",
        ]),
        ("Customer / Community", PURPLE, "Built on top of framework\nDistributed via IPM / OX", [
            "Custom Agents (Lab, Pharmacy, Claims)",
            "Custom MCP Servers (EHR wrappers)",
            "Custom Tools (org-specific APIs)",
            "Custom Skills (local conventions)",
            "SE accelerators (common patterns)",
            "ISV tool packs",
            "Training datasets for skills",
            "Regional HL7 profile knowledge",
        ]),
    ]
    for i, (title, clr, subtitle, items) in enumerate(cols):
        x = Inches(0.8) + i * (col_w + Inches(0.25))
        add_rect(slide, x, Inches(1.5), col_w, col_h, RGBColor(0xFF,0xFF,0xFF), radius=True)
        add_rect(slide, x, Inches(1.5), col_w, Pt(4), clr)
        add_text(slide, x + Inches(0.2), Inches(1.7), col_w - Inches(0.4), Inches(0.4),
                 title, size=15, color=clr, bold=True)
        add_text(slide, x + Inches(0.2), Inches(2.1), col_w - Inches(0.4), Inches(0.5),
                 subtitle, size=10, color=MUTED)
        tf = add_text(slide, x + Inches(0.2), Inches(2.7), col_w - Inches(0.4), Inches(3.5),
                      "", size=12, color=DARK_TEXT)
        for item in items:
            add_bullet(tf, item, size=11, color=DARK_TEXT)

    # --- Slide 8: Catalog & Vector Search ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Vector Catalogs: Semantic Search over IRIS Classes", size=28, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(3), Pt(3), ACCENT)

    # Two catalog cards
    for i, (name, count, scope, kinds, clr) in enumerate([
        ("search_ens", "164 classes", "Business Hosts, Services, Processes, Operations, Adapters",
         "BS:53  BO:38  BP:14  MSG:14  OBA:14  UTL:14  IBA:13  DTL:2  PRD:1  BPL:1", TEAL),
        ("search_hs", "58 classes", "HealthShare transformations, FHIR mappers, SDA helpers, HL7 gateways",
         "DTL:10  SCH:9  BS:8  BP:7  BO:7  GATEWAY:5  API:5  FHIR-DTL:3  CDA:2  HS-MESSAGE:1", GREEN),
    ]):
        y = Inches(1.5) + i * Inches(1.6)
        add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.4), RGBColor(0xFF,0xFF,0xFF), radius=True)
        add_rect(slide, Inches(0.8), y, Inches(0.12), Inches(1.4), clr)
        add_text(slide, Inches(1.2), y + Inches(0.1), Inches(3), Inches(0.35),
                 f"{name}  ({count})", size=16, color=DARK_TEXT, bold=True)
        add_text(slide, Inches(1.2), y + Inches(0.45), Inches(10), Inches(0.35),
                 scope, size=12, color=MID_TEXT)
        add_text(slide, Inches(1.2), y + Inches(0.85), Inches(10), Inches(0.35),
                 kinds, size=10, color=MUTED)

    # Key insight
    add_rect(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.8),
             RGBColor(0xFF,0xF3,0xCD),
             "Key lesson: curated prose (class name + description + key parameters) dramatically outperforms raw class dumps in 384-dim embeddings. Auto-generated accessors and storage boilerplate drown out semantic signal.",
             text_color=DARK_TEXT, text_size=12, radius=True)

    # Screenshot
    add_image_safe(slide, os.path.join(IMG, "11_catalogs.png"),
                   Inches(0.8), Inches(5.9), width=Inches(11.5))

    # --- Slide 9: Performance & Token Optimization ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Performance: Token Optimization Results", size=28, color=DARK_TEXT, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2.5), Pt(3), ACCENT)

    # Before / After table as colored blocks
    metrics = [
        ("Tool catalog tokens",    "~15K", "~10K",   "-33%"),
        ("System prompt tokens",   "~5K",  "~2K",    "-60%"),
        ("Simple query (total)",   "~25K", "~15K",   "-40%"),
        ("Complex task (total)",   "~120K","~50K",   "-58%"),
        ("Simple query latency",   "8-12s","3-5s",   "-60%"),
        ("Complex task latency",   "90s+", "45-60s", "-40%"),
    ]
    header_y = Inches(1.5)
    add_text(slide, Inches(1), header_y, Inches(4.5), Inches(0.4),
             "Metric", size=13, color=MUTED, bold=True)
    add_text(slide, Inches(5.5), header_y, Inches(2.2), Inches(0.4),
             "Before", size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.7), header_y, Inches(2.2), Inches(0.4),
             "After", size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(10), header_y, Inches(2), Inches(0.4),
             "Improvement", size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    for i, (metric, before, after, pct) in enumerate(metrics):
        y = Inches(1.95) + i * Inches(0.55)
        bg = RGBColor(0xF8,0xF9,0xFA) if i % 2 == 0 else RGBColor(0xFF,0xFF,0xFF)
        add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.5), bg)
        add_text(slide, Inches(1), y + Inches(0.08), Inches(4.5), Inches(0.35),
                 metric, size=13, color=DARK_TEXT)
        add_text(slide, Inches(5.5), y + Inches(0.08), Inches(2.2), Inches(0.35),
                 before, size=13, color=RED, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(7.7), y + Inches(0.08), Inches(2.2), Inches(0.35),
                 after, size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(10), y + Inches(0.08), Inches(2), Inches(0.35),
                 pct, size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # Strategies
    tf = add_text(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.4),
                  "Six optimization strategies applied:", size=14, color=DARK_TEXT, bold=True)
    strats = [
        ("ToolFilter policy", "strips 15 framework waste tools before each LLM call"),
        ("Concise tool descriptions", "imperative contracts instead of verbose explanations"),
        ("No markdown in responses", "eliminates rendering bugs in chat UI"),
        ("Monitor token budget", "50K cap per turn prevents runaway costs"),
        ("Multi-turn task decomposition", "research -> propose -> build -> validate -> report"),
        ("Compact system prompt", "< 2K tokens; domain knowledge lives in Skills, not in prompt"),
    ]
    for i, (name, desc) in enumerate(strats):
        col = 0 if i < 3 else 1
        row = i % 3
        x = Inches(0.8) + col * Inches(6)
        y = Inches(5.85) + row * Inches(0.45)
        add_text(slide, x, y, Inches(5.5), Inches(0.4),
                 f"{name}: {desc}", size=11, color=MID_TEXT)

    # --- Slide 10: Roadmap & Next Steps ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Roadmap & Next Steps", size=30, color=WHITE, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2), Pt(3), ACCENT)

    # Three columns: Now / Next / Later
    phases = [
        ("Now (Q3 2026)", TEAL, [
            "Fix Skill %OnNew $ZF marshaling bug",
            "Fix Bedrock tool-result round-trip hang",
            "Ship ToolFilter as framework default",
            "Document %AI.ToolMgr.ExecuteTool query path",
            "Package as IPM for field trial (3 customers)",
        ]),
        ("Next (Q4 2026)", AMBER, [
            "Multi-agent support (Lab, Pharmacy, Claims)",
            "Tool authoring wizard in admin UI",
            "Skill content editor with live preview",
            "FHIR STU3 and CDA Transformation and Mapping Catalog coverage",
            "Open Exchange publication",
        ]),
        ("Later (2027)", PURPLE, [
            "Community marketplace for tools and skills",
            "SE accelerator library (20+ production patterns)",
            "Automated migration assistant (upgrade old DTLs)",
            "Multi-namespace agent isolation",
            "Embedded training mode (learn-by-building)",
        ]),
    ]
    col_w = Inches(3.7)
    col_h = Inches(4.5)
    for i, (title, clr, items) in enumerate(phases):
        x = Inches(0.8) + i * (col_w + Inches(0.25))
        add_rect(slide, x, Inches(1.5), col_w, col_h, PANEL_BG, radius=True)
        add_rect(slide, x, Inches(1.5), col_w, Pt(4), clr)
        add_text(slide, x + Inches(0.2), Inches(1.7), col_w - Inches(0.4), Inches(0.4),
                 title, size=16, color=clr, bold=True)
        tf = add_text(slide, x + Inches(0.2), Inches(2.3), col_w - Inches(0.4), Inches(3.5),
                      "", size=13, color=WHITE)
        for item in items:
            add_bullet(tf, item, size=12, color=WHITE)

    # Bottom: bug count reminder
    add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
             "Blockers to resolve before field trial: 2 framework bugs (Skill %OnNew, Bedrock hang)  |  Built on IRIS for Health 2026.2 with %AI Framework build 162.0",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

    path = os.path.join(DOCS, "Agentic_Health_Interop_Engineering.pptx")
    prs.save(path)
    print(f"  saved {path}")


if __name__ == "__main__":
    print("Building presentations...")
    build_exec_deck()
    build_eng_deck()
    print("\nBoth presentations built.")
